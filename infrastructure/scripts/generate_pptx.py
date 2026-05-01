from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color Palette ─────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x0D, 0x0D)
BG_CARD   = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT    = RGBColor(0x00, 0xC8, 0xFF)
ACCENT2   = RGBColor(0x00, 0xFF, 0xA3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
ORANGE    = RGBColor(0xFF, 0xA5, 0x00)
RED_SOFT  = RGBColor(0xFF, 0x80, 0x80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=BG_DARK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, border_col=None, border_pt=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_col:
        s.line.color.rgb = border_col
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, font="Segoe UI", italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text          = text
    r.font.size     = Pt(size)
    r.font.bold     = bold
    r.font.italic   = italic
    r.font.color.rgb = color
    r.font.name     = font
    return box

def hline(slide, l, t, w, color=ACCENT, h_pt=3):
    rect(slide, l, t, w, Pt(h_pt), color)

def arrow_text(slide, l, t):
    """Draw a → using a textbox (avoids connector API)."""
    tx(slide, "→", l, t, Inches(0.4), Inches(0.5),
       size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

def slide_num(slide, n, total=11):
    tx(slide, f"{n} / {total}",
       SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.42),
       Inches(1.1), Inches(0.35),
       size=9, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

def bullets(slide, items, l, t, w, size=14, gap=Inches(0.42)):
    y = t
    for item in items:
        tx(slide, f"▸  {item}", l, y, w, gap, size=size, color=WHITE)
        y += gap
    return y

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)

tx(s, "🐟  Deep Ocean Fish Market",
   Inches(0.75), Inches(1.5), Inches(10), Inches(1.4),
   size=44, bold=True)

hline(s, Inches(0.75), Inches(2.95), Inches(11.8))

tx(s, "A Progressive Web Application digitizing wholesale fish market operations",
   Inches(0.75), Inches(3.1), Inches(11), Inches(0.75),
   size=17, color=GRAY)

badges = [("Next.js 16", ACCENT), ("Express 5", ACCENT2),
          ("MySQL + Redis", ORANGE), ("PWA ✔", ACCENT)]
bx = Inches(0.75)
for lbl, col in badges:
    rect(s, bx, Inches(4.15), Inches(1.6), Inches(0.45), BG_CARD, col, 1.2)
    tx(s, lbl, bx + Inches(0.08), Inches(4.18), Inches(1.45), Inches(0.4),
       size=11, bold=True, color=col)
    bx += Inches(1.75)

tx(s, "Anthony Hadin  ·  Full-Stack PWA  ·  2025",
   Inches(0.75), Inches(5.2), Inches(9), Inches(0.5), size=13, color=GRAY)
slide_num(s, 1)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Problem Statement", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(4))

problems = [
    "Fish markets rely entirely on paper ledgers — error-prone and slow",
    "No real-time visibility into daily sales, expenses, or profit",
    "Agents & owners are disconnected — no shared live data layer",
    "Offline environments (poor connectivity) are completely unsupported",
    "Manual salary & payroll tracking leads to frequent disputes",
    "No audit trail or PDF invoice generation for buyers",
]
bullets(s, problems, Inches(0.75), Inches(1.3), Inches(11.8), size=16)

rect(s, Inches(0.75), Inches(5.45), Inches(11.8), Inches(0.8), BG_CARD, ACCENT, 1)
tx(s, "🎯  Goal: Replace paper-based workflows with a resilient, offline-first digital system.",
   Inches(0.95), Inches(5.52), Inches(11.4), Inches(0.65),
   size=14, bold=True, color=ACCENT)
slide_num(s, 2)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Tech Stack
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Tech Stack", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(3))

cols = [
    ("Frontend",     ACCENT,  ["Next.js 16  (React 19)", "Tailwind CSS v4", "Framer Motion",
                                "Axios + IndexedDB", "jsPDF  ·  Lucide Icons", "@react-oauth/google"]),
    ("Backend",      ACCENT2, ["Node.js v20 + Express 5", "MySQL 8  (mysql2)", "Redis 7  (ioredis)",
                                "JWT HttpOnly cookies", "Helmet  ·  Rate Limit", "Web Push  ·  PM2"]),
    ("Testing/Ops",  ORANGE,  ["Playwright  (E2E)", "Vitest  ·  Supertest", "Sentry  (observability)",
                                "GitHub Actions CI/CD", "Docker Compose", "Vercel  ·  VPS/PM2"]),
]
cx = Inches(0.75)
for title, col, items in cols:
    rect(s, cx, Inches(1.28), Inches(3.9), Inches(5.7), BG_CARD, col, 1.2)
    tx(s, title, cx + Inches(0.15), Inches(1.38), Inches(3.6), Inches(0.5),
       size=15, bold=True, color=col)
    hline(s, cx + Inches(0.15), Inches(1.88), Inches(3.6), col, 2)
    y = Inches(2.05)
    for item in items:
        tx(s, f"  {item}", cx + Inches(0.15), y, Inches(3.6), Inches(0.42), size=13, color=WHITE)
        y += Inches(0.42)
    cx += Inches(4.27)
slide_num(s, 3)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – System Architecture
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "System Architecture", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(5))

def arch_box(slide, label, sub, l, t, w=Inches(2.5), h=Inches(1.0), col=ACCENT):
    rect(slide, l, t, w, h, BG_CARD, col, 1.5)
    tx(slide, label, l + Inches(0.1), t + Inches(0.07), w - Inches(0.2), Inches(0.44),
       size=12, bold=True, color=col)
    tx(slide, sub,   l + Inches(0.1), t + Inches(0.5),  w - Inches(0.2), Inches(0.45),
       size=10, color=GRAY)

# Layer 1: Client
arch_box(s, "Browser / PWA",  "Service Worker\n+ IndexedDB Cache",
         Inches(0.75), Inches(1.4), col=ACCENT)
arch_box(s, "Next.js Client", "React 19 · Tailwind\nFramer Motion",
         Inches(3.55), Inches(1.4), col=ACCENT)
arrow_text(s, Inches(3.1), Inches(1.65))

# Layer 1: API + Auth
arch_box(s, "Express 5 API",  "JWT Middleware\nHelmet · Rate Limit",
         Inches(6.35), Inches(1.4), col=ACCENT2)
arrow_text(s, Inches(5.9), Inches(1.65))
arch_box(s, "Auth Service",   "Google OAuth\nBcrypt · Denylist",
         Inches(9.15), Inches(1.4), col=ORANGE)
arrow_text(s, Inches(8.7), Inches(1.65))

# Arrow down
tx(s, "↓", Inches(7.45), Inches(2.38), Inches(0.5), Inches(0.45),
   size=20, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
tx(s, "↓", Inches(10.2), Inches(2.38), Inches(0.5), Inches(0.45),
   size=20, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Layer 2: Data
arch_box(s, "MySQL 8",   "Sales · Boats\nUsers · Salaries",
         Inches(5.5), Inches(3.1), col=ACCENT2)
arch_box(s, "Redis 7",   "Token Denylist\nCache · Queues",
         Inches(8.3), Inches(3.1), col=ORANGE)
arch_box(s, "Web Push",  "Push Notifications\nVAPID Keys",
         Inches(11.0), Inches(3.1), w=Inches(2.1), col=DARK_GRAY)

tx(s, "⚡ Offline: mutations queued in IndexedDB → Service Worker syncs on reconnect",
   Inches(0.75), Inches(5.65), Inches(12), Inches(0.55),
   size=12, color=ACCENT2, italic=True)
slide_num(s, 4)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Folder Structure
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Folder Structure Overview", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(5))

tree = (
    "FISH_MARKET/\n"
    "├── client/                   → Next.js 16 PWA Frontend\n"
    "│   ├── src/app/(auth)/        → Login · Register · Forgot password\n"
    "│   ├── src/app/(dashboard)/   → Owner dashboard · Agent dashboard\n"
    "│   ├── src/app/customer/      → Buyer portal\n"
    "│   ├── src/contexts/          → Auth · Language · Notifications\n"
    "│   ├── src/components/        → UI primitives · VoiceInput\n"
    "│   └── public/sw.js           → Service Worker (offline sync)\n"
    "│\n"
    "├── server/                   → Express 5 REST API\n"
    "│   ├── src/routes/            → auth / boats / sales / salaries\n"
    "│   ├── src/middleware/         → JWT auth · rate limiting\n"
    "│   └── src/db/                → MySQL pool · Redis client\n"
    "│\n"
    "├── tests/                    → Playwright E2E suite\n"
    "├── shared/                   → Types shared by client & server\n"
    "├── docker-compose.yml         → MySQL + Redis local environment\n"
    "└── .github/workflows/         → CI/CD (lint → test → build → deploy)"
)
rect(s, Inches(0.75), Inches(1.25), Inches(12.3), Inches(5.8), BG_CARD)
tx(s, tree, Inches(0.9), Inches(1.32), Inches(12.0), Inches(5.65),
   size=11, color=ACCENT2, font="Consolas")
slide_num(s, 5)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Core Features
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Core Features", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(3))

features = [
    ("🌐  PWA + Offline-First",
     "Installable on mobile/desktop. SW caches shell & queues mutations in IndexedDB for sync on reconnect."),
    ("🎙  Voice Commerce UI",
     "Hands-free voice input with Tamil number parsing & Voice Queue for busy market environments."),
    ("🔐  Role-Based Access",
     "3 distinct roles: Owner, Agent, Buyer — isolated dashboards with middleware-guarded routes."),
    ("📊  Live Analytics",
     "Daily/weekly sales trends, boat performance charts, net profit across boats with PDF export."),
    ("💰  Payroll Manager",
     "Track staff salaries, advances & deductions — all subtracted automatically from net profit."),
    ("🔔  Push Notifications",
     "Web Push (VAPID) alerts owners & agents on pending boat-link requests & critical events."),
]

positions = [
    (Inches(0.75),  Inches(1.28)),
    (Inches(4.62),  Inches(1.28)),
    (Inches(8.49),  Inches(1.28)),
    (Inches(0.75),  Inches(3.75)),
    (Inches(4.62),  Inches(3.75)),
    (Inches(8.49),  Inches(3.75)),
]
CW, CH = Inches(3.68), Inches(2.22)
for (lbl, desc), (fx, fy) in zip(features, positions):
    rect(s, fx, fy, CW, CH, BG_CARD, ACCENT, 0.8)
    tx(s, lbl, fx + Inches(0.12), fy + Inches(0.1), CW - Inches(0.24), Inches(0.5),
       size=12, bold=True, color=ACCENT)
    hline(s, fx + Inches(0.12), fy + Inches(0.58), CW - Inches(0.24), DARK_GRAY, 1)
    tx(s, desc, fx + Inches(0.12), fy + Inches(0.62), CW - Inches(0.24), Inches(1.5),
       size=11, color=GRAY)
slide_num(s, 6)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – How It Works
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "How It Works", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(3))

steps = [
    ("1", "Open App",          "PWA loads via SW.\nShell served from\ncache instantly."),
    ("2", "Login / OAuth",     "JWT cookie set.\nRole decoded →\nroute redirect."),
    ("3", "Record Sale",       "Agent enters sale\nvia UI or hands-free\nvoice input."),
    ("4", "API Request",       "Axios POSTs to\nExpress. JWT cookie\nauto-attached."),
    ("5", "Persist to DB",     "MySQL stores data.\nRedis updates\natomic totals."),
    ("6", "Owner Views",       "Dashboard fetches\naggregated reports\nin real time."),
]

bx = Inches(0.65)
for i, (num, title, body) in enumerate(steps):
    rect(s, bx + Inches(0.6), Inches(1.55), Inches(0.67), Inches(0.67), ACCENT)
    tx(s, num, bx + Inches(0.6), Inches(1.57), Inches(0.67), Inches(0.62),
       size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    rect(s, bx, Inches(2.28), Inches(2.0), Inches(2.85), BG_CARD, DARK_GRAY, 0.8)
    tx(s, title, bx + Inches(0.1), Inches(2.38), Inches(1.8), Inches(0.48),
       size=13, bold=True, color=WHITE)
    hline(s, bx + Inches(0.1), Inches(2.84), Inches(1.8), ACCENT, 2)
    tx(s, body, bx + Inches(0.1), Inches(2.95), Inches(1.8), Inches(2.0),
       size=11, color=GRAY)
    if i < len(steps) - 1:
        arrow_text(s, bx + Inches(2.0), Inches(3.35))
    bx += Inches(2.22)

tx(s, "⚡ If OFFLINE at step 3:  mutation saved to IndexedDB → SW syncs automatically on reconnect",
   Inches(0.75), Inches(5.58), Inches(12), Inches(0.65),
   size=13, color=ACCENT2, italic=True)
slide_num(s, 7)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Challenges & Solutions
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Challenges & Solutions", Inches(0.75), Inches(0.38), Inches(10), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(5))

challenges = [
    ("🔄  Offline Sync Race Conditions",
     "Queued mutations could conflict on sync, corrupting daily totals.",
     "Atomic Redis INCRBY counters + idempotency keys make every mutation safe to replay. SW processes IndexedDB queue sequentially."),
    ("🔐  JWT Refresh Loop",
     "Concurrent 401s from multiple tabs triggered multiple token-refresh calls, causing race conditions.",
     "Module-level refreshPromise deduplication: first 401 starts the refresh; all others await the same promise. Retry on success, dispatch auth:unauthorized on failure."),
    ("🎭  Playwright Test Flakiness",
     "E2E tests intermittently failed (429 / 404) due to real network calls and timing issues.",
     "Centralised page.route() API mocking + server health-check polling before tests. Visual regression snapshots with DPR normalisation committed to repo."),
]

ty = Inches(1.3)
for title, problem, solution in challenges:
    rect(s, Inches(0.75), ty, Inches(12.3), Inches(1.72), BG_CARD, DARK_GRAY, 0.8)
    tx(s, title, Inches(0.95), ty + Inches(0.08), Inches(12.0), Inches(0.45),
       size=14, bold=True, color=ACCENT)
    tx(s, f"Problem: {problem}", Inches(0.95), ty + Inches(0.52), Inches(12.0), Inches(0.45),
       size=11, color=RED_SOFT)
    tx(s, f"Solution: {solution}", Inches(0.95), ty + Inches(0.95), Inches(12.0), Inches(0.65),
       size=11, color=ACCENT2)
    ty += Inches(1.9)
slide_num(s, 8)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Demo / Screenshots
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Demo / Screenshots", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(4))

placeholders = [
    ("Login Page",       Inches(0.75),  Inches(1.28)),
    ("Owner Dashboard",  Inches(4.65),  Inches(1.28)),
    ("Agent Dashboard",  Inches(8.55),  Inches(1.28)),
    ("Sales Entry",      Inches(0.75),  Inches(4.05)),
    ("PDF Invoice",      Inches(4.65),  Inches(4.05)),
    ("Mobile PWA View",  Inches(8.55),  Inches(4.05)),
]
PW, PH = Inches(3.62), Inches(2.55)
for label, px, py in placeholders:
    rect(s, px, py, PW, PH, BG_CARD, DARK_GRAY, 1)
    hline(s, px, py, PW, ACCENT, 2)
    tx(s, f"[ {label} ]", px, py + Inches(1.05), PW, Inches(0.5),
       size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)
slide_num(s, 9)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Future Scope
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), SLIDE_H, ACCENT)
tx(s, "Future Scope", Inches(0.75), Inches(0.38), Inches(9), Inches(0.7),
   size=30, bold=True)
hline(s, Inches(0.75), Inches(1.08), Inches(3))

future = [
    ("📱  Native Mobile App",       "React Native/Expo wrapper for iOS & Android distribution."),
    ("🤖  AI Price Prediction",     "LSTM/ARIMA model to forecast daily fish prices from historical data."),
    ("🌐  Multi-Market Support",    "Multi-tenant architecture to serve multiple independent markets."),
    ("💳  UPI Payment Gateway",     "Direct UPI integration so buyers can pay in-app."),
    ("📦  Inventory Management",    "Track stock levels, catch arrivals, and wastage per boat."),
    ("🌍  Multilingual Expansion",  "Extend voice & UI beyond Tamil → Malayalam, Telugu, Kannada."),
]
ly = Inches(1.38)
for title, desc in future:
    rect(s, Inches(0.75), ly, Inches(12.3), Inches(0.75), BG_CARD, DARK_GRAY, 0.6)
    tx(s, title, Inches(0.95), ly + Inches(0.1), Inches(4.5), Inches(0.55),
       size=13, bold=True, color=ACCENT)
    tx(s, desc, Inches(5.3), ly + Inches(0.1), Inches(7.5), Inches(0.55),
       size=13, color=GRAY)
    ly += Inches(0.92)
slide_num(s, 10)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Thank You / Q&A
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ACCENT)
rect(s, Inches(0), SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), ACCENT)

tx(s, "Thank You!", Inches(2), Inches(1.2), Inches(9.33), Inches(1.6),
   size=64, bold=True, align=PP_ALIGN.CENTER)
hline(s, Inches(4.17), Inches(2.88), Inches(5), ACCENT, 4)
tx(s, "Q & A", Inches(2), Inches(3.05), Inches(9.33), Inches(1.0),
   size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tx(s, "Deep Ocean Fish Market  ·  Anthony Hadin  ·  Full-Stack PWA",
   Inches(2), Inches(4.35), Inches(9.33), Inches(0.6),
   size=14, color=GRAY, align=PP_ALIGN.CENTER)
tx(s, "github.com/anthonyhadin-coder/FISH_MARKET",
   Inches(2), Inches(5.1), Inches(9.33), Inches(0.5),
   size=13, color=ACCENT2, align=PP_ALIGN.CENTER)
tx(s, "Built with ❤️  using Next.js · Express · MySQL · Redis · Playwright",
   Inches(1.5), Inches(5.9), Inches(10.33), Inches(0.5),
   size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"c:\coding_program_files\FISH_MARKET\FishMarket_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"    Total slides: {len(prs.slides)}")
